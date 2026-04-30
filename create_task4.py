from pptx import Presentation
from pptx.util import Pt

pptx_path = r"C:\Users\Priya\Downloads\Presentation_Template.pptx"
out_path = r"C:\Users\Priya\.gemini\antigravity\scratch\tata_iq_eda\Task4_AI_Collections_Strategy.pptx"

prs = Presentation(pptx_path)

def set_font(paragraph, bold=False, size=18):
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.bold = bold

# Slide 2: How the System Works
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame and "Inputs:" in shape.text:
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "Inputs: Behavioral financial data (Credit Utilization, Recent Payment History, Debt-to-Income Ratio)."
        p.level = 0
        p2 = text_frame.add_paragraph()
        p2.text = "Decision Logic: A Random Forest predictive model calculates delinquency probability and flags accounts exceeding 80% utilization with a single missed payment."
        p2.level = 0
        p3 = text_frame.add_paragraph()
        p3.text = "Actions: Automated personalized outreach (SMS/Email) offering flexible repayment options before escalation."
        p3.level = 0
        p4 = text_frame.add_paragraph()
        p4.text = "Learning Loop: The system monitors repayment success rates of flagged customers and updates the model monthly to refine risk thresholds."
        p4.level = 0
        
        for para in text_frame.paragraphs:
            set_font(para, size=18)

# Slide 3: Role of Agentic AI
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame and "Explain which parts" in shape.text:
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "Autonomous Operations:"
        p.level = 0
        p2 = text_frame.add_paragraph()
        p2.text = "Risk Scoring: Continuous calculation of delinquency probabilities based on live data."
        p2.level = 1
        p3 = text_frame.add_paragraph()
        p3.text = "Tier 1 Outreach: Sending early intervention reminder emails/SMS."
        p3.level = 1
        p4 = text_frame.add_paragraph()
        p4.text = "Triage: Routing high-risk or complex cases to specialized collections agents."
        p4.level = 1
        
        p5 = text_frame.add_paragraph()
        p5.text = "Human Oversight (Human-in-the-Loop):"
        p5.level = 0
        p6 = text_frame.add_paragraph()
        p6.text = "Escalation Handling: Complex negotiations or hardship plan approvals."
        p6.level = 1
        p7 = text_frame.add_paragraph()
        p7.text = "Model Auditing: Monthly reviews of the model's fairness and accuracy metrics."
        p7.level = 1
        p8 = text_frame.add_paragraph()
        p8.text = "High-Value Account Decisions: Final approval before aggressive collection on large balances."
        p8.level = 1
        
        for para in text_frame.paragraphs:
            if para.level == 0:
                set_font(para, bold=True, size=20)
            else:
                set_font(para, size=18)

# Slide 4: Responsible AI Guardrails
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.has_text_frame and "List and describe" in shape.text:
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "Disparate Impact Monitoring: Continuous fairness checks to ensure no demographic group (Age/Location) is disproportionately penalized."
        p.level = 0
        p2 = text_frame.add_paragraph()
        p2.text = "Explainability (XAI): Model provides feature-level rationale (e.g., 'Flagged due to 90% utilization') for every prediction."
        p2.level = 0
        p3 = text_frame.add_paragraph()
        p3.text = "Data Privacy & Compliance: Strictly adheres to ECOA by excluding protected demographic variables from decision logic."
        p3.level = 0
        p4 = text_frame.add_paragraph()
        p4.text = "Human-in-the-Loop: Critical escalation pathways mandate human agent review, ensuring AI acts as an assistant rather than final arbiter for severe actions."
        p4.level = 0
        
        for para in text_frame.paragraphs:
            set_font(para, size=18)

# Slide 5: Expected Business Impact
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame and "Summarize in 1-2 slides" in shape.text:
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "Business KPIs:"
        p.level = 0
        p2 = text_frame.add_paragraph()
        p2.text = "15% reduction in 60-day delinquency rates via early intervention."
        p2.level = 1
        p3 = text_frame.add_paragraph()
        p3.text = "Decreased operational costs by automating Tier 1 outreach."
        p3.level = 1
        p4 = text_frame.add_paragraph()
        p4.text = "Increased total debt recovery percentage across the portfolio."
        p4.level = 1
        
        p5 = text_frame.add_paragraph()
        p5.text = "Customer Outcomes:"
        p5.level = 0
        p6 = text_frame.add_paragraph()
        p6.text = "Improved customer trust through supportive, non-punitive initial contact."
        p6.level = 1
        p7 = text_frame.add_paragraph()
        p7.text = "Tailored flexible repayment options preventing credit score damage."
        p7.level = 1
        p8 = text_frame.add_paragraph()
        p8.text = "Fairer assessments completely detached from demographic biases."
        p8.level = 1
        
        for para in text_frame.paragraphs:
            if para.level == 0:
                set_font(para, bold=True, size=20)
            else:
                set_font(para, size=18)

# Remove Slide 6 if present
if len(prs.slides) > 5:
    rId = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]

prs.save(out_path)
print("Saved pptx to:", out_path)
