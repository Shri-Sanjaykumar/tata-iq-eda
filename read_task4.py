from pptx import Presentation

pptx_path = r"C:\Users\Priya\Downloads\Presentation_Template.pptx"

try:
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
except Exception as e:
    print(f"Error reading presentation: {e}")
